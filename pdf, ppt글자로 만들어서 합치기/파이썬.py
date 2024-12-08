import os
from PyPDF2 import PdfReader
from pptx import Presentation
from tkinter import Tk, filedialog

def convert_pdf_to_text(pdf_path):
    text = ""
    try:
        reader = PdfReader(pdf_path)
        for page in reader.pages:
            text += page.extract_text()
    except Exception as e:
        text = f"Error processing {pdf_path}: {e}"
    return text

def convert_ppt_to_text(ppt_path):
    text = ""
    try:
        presentation = Presentation(ppt_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        text = f"Error processing {ppt_path}: {e}"
    return text

def process_files(file_paths, output_file):
    with open(output_file, 'w', encoding='utf-8') as output:
        for file_path in file_paths:
            filename = os.path.basename(file_path)
            if file_path.lower().endswith('.pdf'):
                file_text = convert_pdf_to_text(file_path)
            elif file_path.lower().endswith('.ppt') or file_path.lower().endswith('.pptx'):
                file_text = convert_ppt_to_text(file_path)
            else:
                continue
            
            output.write(f"{filename}\n")
            output.write(f"{file_text}\n")
            output.write("\n")

def main():
    # Tkinter 초기화
    root = Tk()
    root.withdraw()  # GUI 창 숨기기

    # 파일 선택 창 열기
    file_paths = filedialog.askopenfilenames(
        title="Select PDF and PPT files",
        filetypes=[("PDF files", "*.pdf"), ("PPT files", "*.ppt;*.pptx")],
    )
    
    if not file_paths:
        print("No files selected.")
        return

    # 출력 파일 경로 설정
    output_file_path = filedialog.asksaveasfilename(
        title="Save combined text file",
        defaultextension=".txt",
        filetypes=[("Text files", "*.txt")]
    )

    if not output_file_path:
        print("No output file selected.")
        return

    # 파일을 처리하여 합치기
    process_files(file_paths, output_file_path)
    print(f"Combined text file saved to {output_file_path}")

if __name__ == "__main__":
    main()
