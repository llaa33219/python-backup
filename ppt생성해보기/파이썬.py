from pptx import Presentation
from pptx.util import Inches

# 새 프레젠테이션 생성
prs = Presentation()

# 슬라이드 1: 프레젠테이션 제목
slide_title = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_title.shapes.title
subtitle = slide_title.placeholders[1]

title.text = "ChatGPT Team Plan"
subtitle.text = "An overview of our strategic vision and objectives."

# 슬라이드 2: 목표 소개
slide_goals = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide_goals.shapes.title, slide_goals.placeholders[1]

title.text = "Our Goals"
content.text = "• Enhance ChatGPT's capabilities\n• Expand user engagement\n• Develop new features"

# 슬라이드 3: 팀 구성
slide_team = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide_team.shapes.title, slide_team.placeholders[1]

title.text = "Team Structure"
content.text = "• Research and Development Team\n• Marketing and Outreach Team\n• Customer Support Team"

# 파일 저장
prs.save('ChatGPT_Team_Plan_Presentation.pptx')
